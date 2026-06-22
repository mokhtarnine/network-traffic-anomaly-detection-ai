#!/usr/bin/env python3
"""
PRESENTATION_final.pptx  —  One accent (amber), no code, feature-focused,
architecture diagram, limitations, annotated screenshots.
"""
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

# ── Palette — ONE accent (amber), everything else is grey scale ──────
BG    = RGBColor(0x0c, 0x0f, 0x1a)   # near-black
SURF  = RGBColor(0x14, 0x1c, 0x2e)   # card surface
SURF2 = RGBColor(0x1a, 0x24, 0x38)   # slightly lighter card
LINE  = RGBColor(0x1e, 0x2d, 0x42)   # separators / borders
AMB   = RGBColor(0xe2, 0xa8, 0x00)   # amber — THE accent (used sparingly)
AMB2  = RGBColor(0x78, 0x56, 0x00)   # dim amber for decorative bg numbers
W1    = RGBColor(0xf8, 0xfa, 0xfc)   # primary text
W2    = RGBColor(0x94, 0xa3, 0xb8)   # secondary text
W3    = RGBColor(0x47, 0x56, 0x69)   # muted / labels

SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.85)
RM = Inches(0.85)
TM = Inches(0.58)

SS = Path(__file__).parent / "screenshots"

# ── Primitives ────────────────────────────────────────────────────────

def prs_new():
    p = Presentation()
    p.slide_width  = SW
    p.slide_height = SH
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(sl, l, t, w, h, fill, border=None, bpt=0.8):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bpt)
    else:
        s.line.color.rgb = fill; s.line.width = Pt(0)
    return s

def txt(sl, text, l, t, w, h, sz, color=W1, bold=False,
        align=PP_ALIGN.LEFT, italic=False, name="Calibri"):
    if not text: return
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    pg = tf.paragraphs[0]; pg.alignment = align
    r  = pg.add_run()
    r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = name

def rule(sl, y, l=None, rw=None, c=LINE, th=1.0):
    l  = l  if l  is not None else LM
    rw = rw if rw is not None else (SW - LM - RM)
    box(sl, l, y, rw, Inches(th / 72), fill=c)

def bg(sl):
    box(sl, 0, 0, SW, SH, fill=BG)

# Decorative huge background number — gives creative texture per slide
def bignumber_bg(sl, num_str):
    txt(sl, num_str,
        SW - Inches(3.8), Inches(-0.2), Inches(3.8), Inches(2.8),
        144, AMB2, True, PP_ALIGN.RIGHT, name="Calibri")

def header(sl, title, subtitle=None):
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)          # left amber bar
    txt(sl, title, LM, TM, SW - LM - RM, Inches(0.5), 26, W1, True)
    if subtitle:
        txt(sl, subtitle, LM, TM + Inches(0.46),
            SW - LM - RM, Inches(0.28), 11.5, W2)
    rule(sl, TM + (Inches(0.82) if subtitle else Inches(0.56)))

def content_y(has_subtitle=True):
    return TM + (Inches(0.96) if has_subtitle else Inches(0.72))

def pgnum(sl, n, total=14):
    txt(sl, f"{n}  /  {total}", SW - RM - Inches(1.1), SH - Inches(0.42),
        Inches(1.1), Inches(0.3), 9, W3, align=PP_ALIGN.RIGHT)

# ── Screenshot + annotation engine ───────────────────────────────────
# Screenshots: 1400 × 850 px
# Image placed: left=LM, top=0.9in, width=SW-2*LM, height=SH-0.9-0.15
IL = LM
IT = Inches(0.88)
IW = SW - 2 * LM
IH = SH - IT - Inches(0.14)
IPW, IPH = 1400, 850

def px_to(x): return IL + (x / IPW) * IW
def py_to(y): return IT + (y / IPH) * IH

def callout(sl, sx, sy, label, direction="right",
            line_len=Inches(0.5), lw=Inches(2.5)):
    """Dot + line + label box annotation."""
    D  = Inches(0.09)
    LH = Inches(0.26)
    TK = Inches(0.013)

    box(sl, sx - D/2, sy - D/2, D, D, fill=AMB)    # dot

    if direction == "right":
        box(sl, sx + D/2, sy - TK/2, line_len, TK, fill=AMB)
        bx, by = sx + D/2 + line_len, sy - LH/2
    elif direction == "left":
        box(sl, sx - D/2 - line_len, sy - TK/2, line_len, TK, fill=AMB)
        bx, by = sx - D/2 - line_len - lw, sy - LH/2
    elif direction == "up":
        box(sl, sx - TK/2, sy - D/2 - line_len, TK, line_len, fill=AMB)
        bx, by = sx - lw/2, sy - D/2 - line_len - LH
    elif direction == "down":
        box(sl, sx - TK/2, sy + D/2, TK, line_len, fill=AMB)
        bx, by = sx - lw/2, sy + D/2 + line_len

    box(sl, bx, by, lw, LH, fill=SURF2, border=AMB, bpt=0.9)
    txt(sl, label, bx + Inches(0.1), by + Inches(0.02),
        lw - Inches(0.12), LH, 8.5, AMB, True)

def ss_slide(prs, title, subtitle, fname):
    """Screenshot slide with dark header bar."""
    sl = blank(prs); bg(sl)
    hh = Inches(0.82)
    box(sl, 0, 0, SW, hh, fill=SURF)
    box(sl, 0, hh, SW, Inches(0.04), fill=AMB)
    box(sl, 0, 0, Inches(0.1), hh, fill=AMB)
    txt(sl, title,    Inches(0.24), Inches(0.08), SW - Inches(1.5), Inches(0.38), 17, W1, True)
    txt(sl, subtitle, Inches(0.24), Inches(0.50), SW - Inches(0.4), Inches(0.26), 10, W2)
    img = SS / fname
    if img.exists():
        sl.shapes.add_picture(str(img), IL, IT, IW, IH)
    return sl

# =======================================================================
# SLIDES
# =======================================================================

# ── 1 · TITLE ────────────────────────────────────────────────────────
def s01(prs):
    sl = blank(prs); bg(sl)

    # Full-height amber left bar
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)

    # Decorative right panel
    box(sl, SW - Inches(4.2), 0, Inches(4.2), SH, fill=SURF)
    box(sl, SW - Inches(4.2), 0, Inches(0.04), SH, fill=LINE)

    # Institution
    txt(sl, "SUPMTI Rabat   AI Capstone 2025-2026   Topic 15",
        LM, Inches(1.1), Inches(8.5), Inches(0.3), 10, W3)

    # Main title
    txt(sl, "Network Traffic", LM, Inches(1.6), Inches(9), Inches(1.05), 58, W1, True)
    txt(sl, "Anomaly Detection", LM, Inches(2.55), Inches(9), Inches(1.05), 58, AMB, True)

    # Tagline
    txt(sl, "Unsupervised K-Means Clustering   NSL-KDD Dataset",
        LM, Inches(3.72), Inches(9), Inches(0.36), 13, W2)

    rule(sl, Inches(4.22), LM, Inches(4.2), LINE)

    # Author block
    for i, (lbl, val) in enumerate([
        ("Author",     "Said Louham"),
        ("Supervisor", "Pr. Soufiane HAMIDA"),
        ("Year",       "2025 / 2026"),
    ]):
        x = LM + i * Inches(3.4)
        txt(sl, lbl.upper(), x, Inches(4.4), Inches(3.2), Inches(0.26), 7.5, W3, True)
        txt(sl, val,          x, Inches(4.64), Inches(3.2), Inches(0.36), 13.5, W1, True)

    # Right panel content
    rx = SW - Inches(4.0)
    txt(sl, "Built with", rx, Inches(1.1), Inches(3.6), Inches(0.28), 9, W3)
    for i, tech in enumerate(["Python 3.12", "Streamlit", "scikit-learn",
                               "Plotly", "SQLite", "pytest"]):
        ty = Inches(1.5) + i * Inches(0.44)
        box(sl, rx, ty + Inches(0.1), Inches(0.06), Inches(0.2), fill=AMB)
        txt(sl, tech, rx + Inches(0.18), ty, Inches(3.4), Inches(0.36), 12, W1, True)

# ── 2 · THE CHALLENGE ────────────────────────────────────────────────
def s02(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "01")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "The Challenge", "Why automated anomaly detection matters")
    pgnum(sl, 1)
    y0 = content_y()

    rows = [
        ("Volume",
         "Networks generate thousands of connections per second. "
         "A single 1Gbps link can produce over 60,000 flows per minute — "
         "manual inspection is physically impossible."),
        ("Variety",
         "Attacks differ fundamentally in behavior. "
         "DoS floods bandwidth. Probes scan for open ports quietly. "
         "R2L and U2R exploit service vulnerabilities through legitimate channels. "
         "No single rule catches all four."),
        ("No Labels",
         "In production, traffic arrives without pre-labeled categories. "
         "Supervised classifiers require labeled attack samples at training time — "
         "a requirement that cannot always be met in new environments."),
    ]
    rh = (SH - y0 - Inches(1.1)) / 3
    for i, (title, desc) in enumerate(rows):
        y = y0 + i * rh
        box(sl, LM, y + Inches(0.06), SW - LM - RM, rh - Inches(0.08), fill=SURF)
        box(sl, LM, y + Inches(0.06), Inches(0.06), rh - Inches(0.08), fill=AMB)
        txt(sl, f"0{i+1}", LM + Inches(0.2), y + Inches(0.2),
            Inches(0.55), Inches(0.58), 28, AMB, True)
        txt(sl, title, LM + Inches(0.88), y + Inches(0.14),
            SW - LM - RM - Inches(1.0), Inches(0.34), 15, W1, True)
        txt(sl, desc, LM + Inches(0.88), y + Inches(0.5),
            SW - LM - RM - Inches(1.0), rh - Inches(0.6), 11, W2)

    # Solution strip
    sy = y0 + 3 * rh + Inches(0.08)
    box(sl, LM, sy, SW - LM - RM, Inches(0.54), fill=SURF2, border=AMB)
    txt(sl, "Solution  —  K-Means unsupervised clustering groups traffic by behavior. "
            "No labels required at training time. "
            "The cluster with the highest attack density is automatically flagged.",
        LM + Inches(0.22), sy + Inches(0.1),
        SW - LM - RM - Inches(0.28), Inches(0.36), 11.5, W1)

# ── 3 · DATASET ──────────────────────────────────────────────────────
def s03(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "02")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "NSL-KDD Dataset", "Improved KDD Cup 1999 — standard academic benchmark for intrusion detection research")
    pgnum(sl, 2)
    y0 = content_y()

    sw4 = (SW - LM - RM - 3 * Inches(0.14)) / 4
    sh4 = Inches(1.88)
    for i, (val, lbl, sub) in enumerate([
        ("125,973", "Training records", "KDDTrain+.txt"),
        ("22,544",  "Test records",     "KDDTest+.txt"),
        ("41",      "Features/record",  "per connection"),
        ("23",      "Attack classes",   "in training set"),
    ]):
        x = LM + i * (sw4 + Inches(0.14))
        box(sl, x, y0, sw4, sh4, fill=SURF)
        box(sl, x, y0, sw4, Inches(0.05), fill=AMB)
        txt(sl, val, x, y0 + Inches(0.14), sw4, Inches(0.95), 42, AMB, True, PP_ALIGN.CENTER)
        txt(sl, lbl, x, y0 + Inches(1.06), sw4, Inches(0.3),  12, W1,  True, PP_ALIGN.CENTER)
        txt(sl, sub, x, y0 + Inches(1.36), sw4, Inches(0.26), 9.5, W3, False, PP_ALIGN.CENTER)

    rule(sl, y0 + sh4 + Inches(0.22))
    y2 = y0 + sh4 + Inches(0.38)
    cw = (SW - LM - RM - Inches(0.3)) / 2
    rx = LM + cw + Inches(0.3)

    txt(sl, "Attack Categories", LM, y2, cw, Inches(0.3), 12, AMB, True)
    attacks = [
        ("DoS",    "neptune, smurf, back, pod, teardrop"),
        ("Probe",  "ipsweep, nmap, portsweep, satan"),
        ("R2L",    "warezclient, imap, spy, ftp_write"),
        ("U2R",    "buffer_overflow, rootkit, perl"),
        ("Normal", "approx. 55% of training records"),
    ]
    for i, (cat, ex) in enumerate(attacks):
        yy = y2 + Inches(0.36) + i * Inches(0.43)
        box(sl, LM, yy + Inches(0.09), Inches(0.06), Inches(0.24),
            fill=AMB if cat != "Normal" else W3)
        txt(sl, cat, LM + Inches(0.18), yy, Inches(1.1), Inches(0.36), 11, W1, True)
        txt(sl, ex,  LM + Inches(1.32), yy, cw - Inches(1.4), Inches(0.36), 10.5, W2)

    txt(sl, "Dataset Properties", rx, y2, cw, Inches(0.3), 12, AMB, True)
    for i, f in enumerate([
        "No missing values in either file",
        "Pre-split train / test — zero data leakage",
        "3 categorical  +  38 numeric features",
        "Labels removed before training",
        "Labels used only after clustering to evaluate quality",
    ]):
        yy = y2 + Inches(0.36) + i * Inches(0.43)
        box(sl, rx, yy + Inches(0.1), Inches(0.06), Inches(0.22), fill=W3)
        txt(sl, f, rx + Inches(0.18), yy, cw - Inches(0.22), Inches(0.36), 10.5, W2)

# ── 4 · ARCHITECTURE ─────────────────────────────────────────────────
def s04(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "03")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "System Architecture", "How the modules connect — from raw file to export")
    pgnum(sl, 3)
    y0 = content_y()

    content_h = SH - y0 - Inches(0.3)

    # ── Layer labels (left column) ───────────────────────
    lbl_w = Inches(1.55)
    for label, cy, ch in [
        ("Interface",  y0,                  Inches(1.0)),
        ("Logic",      y0 + Inches(1.22),   Inches(1.88)),
        ("Storage",    y0 + Inches(3.32),   Inches(1.2)),
    ]:
        box(sl, LM, cy + Inches(0.08), lbl_w, ch - Inches(0.12), fill=SURF)
        box(sl, LM, cy + Inches(0.08), Inches(0.06), ch - Inches(0.12), fill=AMB)
        txt(sl, label, LM + Inches(0.18), cy + ch/2 - Inches(0.2),
            lbl_w - Inches(0.24), Inches(0.36), 11, W2, True)

    content_x = LM + lbl_w + Inches(0.2)
    content_w  = SW - content_x - RM

    # ── Row 1: Streamlit App bar ─────────────────────────
    app_y = y0; app_h = Inches(0.95)
    box(sl, content_x, app_y, content_w, app_h, fill=SURF2, border=AMB, bpt=1.0)
    txt(sl, "Streamlit Application   —   src/app/app.py",
        content_x + Inches(0.2), app_y + Inches(0.14),
        content_w - Inches(0.3), Inches(0.32), 13, W1, True)
    txt(sl, "6 render functions   Session state management   Custom CSS   Sidebar settings panel",
        content_x + Inches(0.2), app_y + Inches(0.5),
        content_w - Inches(0.3), Inches(0.28), 9.5, W2)

    # Connector line down
    cx_mid = content_x + content_w / 2
    box(sl, cx_mid - Inches(0.01), app_y + app_h, Inches(0.02), Inches(0.2), fill=LINE)

    # ── Row 2: 5 module boxes ────────────────────────────
    mods = [
        ("loader.py",        "File loading\nSchema validation"),
        ("eda_analysis.py",  "5 Plotly chart\nfunctions"),
        ("preprocessor.py",  "OHE + Scaler\n43 -> 122 features"),
        ("kmeans_model.py",  "Train, predict\nAnomaly mask"),
        ("metrics.py",       "Silhouette, DB\nARI, P, R, F1"),
    ]
    mod_y = app_y + app_h + Inches(0.2)
    mod_h = Inches(1.7)
    mod_w = (content_w - 4 * Inches(0.12)) / 5
    for i, (name, desc) in enumerate(mods):
        mx = content_x + i * (mod_w + Inches(0.12))
        box(sl, mx, mod_y, mod_w, mod_h, fill=SURF)
        box(sl, mx, mod_y, mod_w, Inches(0.04), fill=AMB)
        txt(sl, name, mx + Inches(0.1), mod_y + Inches(0.1),
            mod_w - Inches(0.14), Inches(0.32), 10.5, AMB, True, PP_ALIGN.CENTER)
        rule(sl, mod_y + Inches(0.46), mx + Inches(0.12), mod_w - Inches(0.24), LINE)
        txt(sl, desc, mx + Inches(0.1), mod_y + Inches(0.56),
            mod_w - Inches(0.14), mod_h - Inches(0.66), 10, W2, False, PP_ALIGN.CENTER)
        # connector down
        box(sl, mx + mod_w/2 - Inches(0.01), mod_y + mod_h,
            Inches(0.02), Inches(0.2), fill=LINE)

    # ── Row 3: 4 storage items ───────────────────────────
    stor_y = mod_y + mod_h + Inches(0.2)
    stor_h = Inches(1.08)
    stores = [
        ("data/raw/",         "KDDTrain+.txt\nKDDTest+.txt",  0, 2),   # spans cols 0-1
        ("models/",           "kmeans_model.pkl\npreprocessor.pkl", 2, 1),
        ("database.py",       "SQLite history\nauto-migration", 3, 1),
        ("PDF / CSV",         "fpdf2 report\nDownload button", 4, 1),
    ]
    for (name, desc, col, span) in stores:
        sx = content_x + col * (mod_w + Inches(0.12))
        sw_s = span * mod_w + (span - 1) * Inches(0.12)
        box(sl, sx, stor_y, sw_s, stor_h, fill=SURF)
        box(sl, sx, stor_y + stor_h - Inches(0.04), sw_s, Inches(0.04), fill=LINE)
        txt(sl, name, sx + Inches(0.12), stor_y + Inches(0.08),
            sw_s - Inches(0.18), Inches(0.28), 10, W1, True)
        txt(sl, desc, sx + Inches(0.12), stor_y + Inches(0.36),
            sw_s - Inches(0.18), stor_h - Inches(0.44), 9.5, W3)

    # ── database.py link to app ──────────────────────────
    box(sl, content_x + Inches(0.2), y0 + Inches(0.7),
        Inches(0.8), Inches(0.02), fill=LINE)

# ── 5 · PIPELINE ─────────────────────────────────────────────────────
def s05(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "04")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "6-Step Pipeline", "End-to-end workflow — every step is automated and interactive")
    pgnum(sl, 4)
    y0 = content_y()

    steps = [
        ("01", "Upload & Validate",    "CSV, Excel, or NSL-KDD .txt\nAuto schema check  |  Row/null/duplicate count"),
        ("02", "Explore",              "4 Plotly chart tabs: Labels, Categoricals, Numerics, Correlation\nAll charts update when new data is uploaded"),
        ("03", "Preprocess",           "OneHotEncoder  +  StandardScaler via ColumnTransformer\n43 raw columns  -->  122 processed features"),
        ("04", "Detect",               "K-Means with configurable k (2-20) via sidebar slider\nPre-trained k=5 model loads in one click"),
        ("05", "Evaluate",             "Silhouette Score  |  Davies-Bouldin  |  ARI  |  Precision  |  Recall  |  F1\nPCA 2D scatter + cluster bar chart + error breakdown"),
        ("06", "Export",               "Download CSV with cluster + is_anomaly columns\nAuto-generated PDF report  |  SQLite history database"),
    ]

    bw  = (SW - LM - RM - 5 * Inches(0.12)) / 6
    bh  = SH - y0 - Inches(0.4)
    for i, (num, title, desc) in enumerate(steps):
        x   = LM + i * (bw + Inches(0.12))
        box(sl, x, y0, bw, bh, fill=SURF)
        box(sl, x, y0, bw, Inches(0.04), fill=AMB)
        txt(sl, num,   x, y0 + Inches(0.1), bw, Inches(0.55), 30, AMB, True, PP_ALIGN.CENTER)
        rule(sl, y0 + Inches(0.68), x + Inches(0.12), bw - Inches(0.24), LINE)
        txt(sl, title, x + Inches(0.1), y0 + Inches(0.76),
            bw - Inches(0.14), Inches(0.38), 13, W1, True, PP_ALIGN.CENTER)
        txt(sl, desc,  x + Inches(0.1), y0 + Inches(1.2),
            bw - Inches(0.14), bh - Inches(1.28), 10, W2, False, PP_ALIGN.CENTER)
        if i < 5:
            txt(sl, ">", x + bw + Inches(0.03), y0 + bh/2 - Inches(0.18),
                Inches(0.12), Inches(0.3), 16, W3, True, PP_ALIGN.CENTER)

# ── 6 · DEMO: UPLOAD ─────────────────────────────────────────────────
def s06(prs):
    sl = ss_slide(prs,
        "Step 1  —  Upload & Validation",
        "Drag-and-drop a file and the app validates schema, counts rows, nulls, and duplicates automatically",
        "02_after_upload.png")
    pgnum(sl, 5)

    callout(sl, px_to(417), py_to(248), "6-step progress indicator", "up",    Inches(0.35), Inches(2.5))
    callout(sl, px_to(465), py_to(300), "File badge with size",       "right", Inches(0.4),  Inches(2.1))
    callout(sl, px_to(840), py_to(385), "Automatic schema validation","right", Inches(0.4),  Inches(2.6))
    callout(sl, px_to(415), py_to(545), "22,544 rows loaded",         "down",  Inches(0.38), Inches(1.9))
    callout(sl, px_to(642), py_to(545), "43 columns",                 "down",  Inches(0.38), Inches(1.6))
    callout(sl, px_to(864), py_to(545), "0 null values",              "down",  Inches(0.38), Inches(1.6))

# ── 7 · DEMO: EDA ────────────────────────────────────────────────────
def s07(prs):
    sl = ss_slide(prs,
        "Step 2  —  Exploratory Data Analysis",
        "4 interactive Plotly tabs — Labels, Categoricals, Numerics, Correlation — rendered on upload",
        "04_eda_labels.png")
    pgnum(sl, 6)

    callout(sl, px_to(440), py_to(488), "4 chart tabs",                   "left",  Inches(0.4),  Inches(1.7))
    callout(sl, px_to(505), py_to(620), "Traffic label distribution",     "left",  Inches(0.4),  Inches(2.3))
    callout(sl, px_to(1060), py_to(600), "Normal vs Attack donut",        "right", Inches(0.38), Inches(2.4))
    callout(sl, px_to(1295), py_to(718), "56.9% of records are attacks",  "right", Inches(0.38), Inches(2.6))

# ── 8 · DEMO: RESULTS ────────────────────────────────────────────────
def s08(prs):
    sl = ss_slide(prs,
        "Step 5  —  Detection Results",
        "Cluster 4 auto-flagged as anomaly cluster (99.6% attack purity)  |  Silhouette Score 0.3930  |  2,219 anomalies",
        "11_cluster_view.png")
    pgnum(sl, 7)

    callout(sl, px_to(724),  py_to(265), "2,219 anomalies detected",          "up",    Inches(0.38), Inches(2.2))
    callout(sl, px_to(1049), py_to(265), "9.84% anomaly rate",                "up",    Inches(0.38), Inches(2.0))
    callout(sl, px_to(850),  py_to(357), "Auto-identified anomaly cluster",   "right", Inches(0.38), Inches(2.8))
    callout(sl, px_to(429),  py_to(501), "Silhouette Score",                  "down",  Inches(0.4),  Inches(1.8))
    callout(sl, px_to(744),  py_to(501), "Davies-Bouldin Index",              "down",  Inches(0.4),  Inches(2.1))
    callout(sl, px_to(530),  py_to(720), "Records per cluster bar chart",     "left",  Inches(0.4),  Inches(2.7))

# ── 9 · DEMO: ERROR ANALYSIS ─────────────────────────────────────────
def s09(prs):
    sl = ss_slide(prs,
        "Step 5  —  Error Analysis",
        "Precision 0.9964  |  Recall 0.1723  |  F1 0.2938  |  10,622 false negatives",
        "12_error_analysis.png")
    pgnum(sl, 8)

    callout(sl, px_to(408),  py_to(501), "Precision 0.9964 — very pure cluster",  "down",  Inches(0.38), Inches(3.0))
    callout(sl, px_to(636),  py_to(501), "Recall 0.1723 — missed mixed attacks",  "down",  Inches(0.38), Inches(3.1))
    callout(sl, px_to(1136), py_to(501), "10,622 false negatives",                "up",    Inches(0.38), Inches(2.0))
    callout(sl, px_to(490),  py_to(655), "neptune = most missed attack type",     "left",  Inches(0.4),  Inches(2.6))

    # WHY explanation strip at bottom over the image
    ey = SH - Inches(1.0)
    box(sl, LM, ey, SW - LM - RM, Inches(0.86), fill=SURF2, border=AMB, bpt=1.0)
    txt(sl, "Why is recall low?",
        LM + Inches(0.22), ey + Inches(0.06), SW - LM - RM - Inches(0.3), Inches(0.28), 11, AMB, True)
    txt(sl, "K-Means distributes all 5 attack types across 5 clusters. "
            "Only ONE cluster is flagged as anomalous. "
            "DoS attacks like neptune that end up in 'normal' clusters are not detected — "
            "this is the inherent trade-off of single-cluster anomaly flagging.",
        LM + Inches(0.22), ey + Inches(0.36), SW - LM - RM - Inches(0.3), Inches(0.44),
        10, W2)

# ── 10 · MODEL SELECTION ─────────────────────────────────────────────
def s10(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "09")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "Model Selection", "Comparing K-Means at k = 3, 4, 5 — and why K-Means was chosen")
    pgnum(sl, 9)
    y0 = content_y()
    lw = Inches(6.6)
    rw = SW - LM - RM - lw - Inches(0.4)
    rx = LM + lw + Inches(0.4)
    bh = SH - y0 - Inches(0.35)

    # Left: table
    trh = Inches(0.56)
    cws = [Inches(1.3), Inches(2.65), Inches(2.65)]
    txt(sl, "Internal clustering metrics (no labels required):",
        LM, y0, lw, Inches(0.28), 10, W3)
    ty = y0 + Inches(0.35)

    for j, h in enumerate(["k", "Silhouette Score  (higher = better)", "Davies-Bouldin  (lower = better)"]):
        x0 = LM + sum(cws[:j])
        box(sl, x0, ty, cws[j], trh, fill=SURF, border=LINE)
        txt(sl, h, x0 + Inches(0.12), ty + Inches(0.16), cws[j] - Inches(0.18), Inches(0.3), 10, AMB, True)

    rows = [("3", "0.4203", "1.0628", False),
            ("4", "0.4254", "0.9609", False),
            ("5", "0.4404", "0.8884", True)]
    for ri, (k, sil, db, win) in enumerate(rows):
        for j, v in enumerate([k, sil, db]):
            x0 = LM + sum(cws[:j])
            yy = ty + (ri + 1) * trh
            bg2 = RGBColor(0x18, 0x26, 0x10) if win else BG
            bc2 = AMB if win else LINE
            box(sl, x0, yy, cws[j], trh, fill=bg2, border=bc2)
            pf  = "* " if win and v == k else ""
            txt(sl, pf + v, x0 + Inches(0.12), yy + Inches(0.16),
                cws[j] - Inches(0.18), Inches(0.3), 11, AMB if win else W1, win)

    wy = ty + 4 * trh + Inches(0.22)
    box(sl, LM, wy, lw, Inches(0.5), fill=RGBColor(0x18, 0x26, 0x10), border=AMB)
    txt(sl, "k = 5 selected — best Silhouette AND lowest Davies-Bouldin simultaneously.",
        LM + Inches(0.18), wy + Inches(0.12), lw - Inches(0.24), Inches(0.3), 11, AMB, True)

    for i, line in enumerate([
        "Why K-Means?",
        "Simple, fast, and interpretable — well-suited for initial anomaly triage on high-dimensional tabular data.",
        "",
        "Why k = 5?",
        "NSL-KDD traffic naturally groups into 5 behaviors: Normal, DoS, Probe, R2L, U2R. Five clusters match this structure without over-fragmenting.",
        "",
        "Why not other algorithms?",
        "DBSCAN requires careful epsilon tuning. Isolation Forest is a stronger anomaly detector but harder to explain to non-technical stakeholders. K-Means produces labeled clusters that are directly interpretable.",
    ]):
        y = wy + Inches(0.7) + i * Inches(0.42)
        if y > SH - Inches(0.5): break
        c = AMB if line in ("Why K-Means?", "Why k = 5?", "Why not other algorithms?") else W2
        b = c == AMB
        if line:
            txt(sl, line, LM, y, lw, Inches(0.38), 10.5 if b else 10, c, b)

    # Right: metric explanations
    box(sl, rx, y0, rw, bh, fill=SURF)
    box(sl, rx, y0, Inches(0.06), bh, fill=AMB)
    txt(sl, "What the metrics measure", rx + Inches(0.2), y0 + Inches(0.16),
        rw - Inches(0.26), Inches(0.3), 12, W1, True)
    rule(sl, y0 + Inches(0.56), rx + Inches(0.2), rw - Inches(0.35))

    explanations = [
        ("Silhouette Score",
         "Measures how tightly packed each cluster is compared to its nearest neighbor cluster. "
         "A score of 0.44 indicates clearly separated clusters with limited overlap."),
        ("Davies-Bouldin Index",
         "Average ratio of intra-cluster spread to inter-cluster distance. "
         "A score of 0.89 (below 1.0) confirms the five clusters are well separated in the 122-dimensional space."),
        ("Adjusted Rand Index",
         "Compares the discovered clusters against the true binary labels. "
         "A perfect match = 1.0. A value of 0.16 reflects the single-cluster limitation — "
         "only one of five attack groups is fully isolated."),
    ]
    ey2 = y0 + Inches(0.66)
    for title, desc in explanations:
        txt(sl, title, rx + Inches(0.22), ey2, rw - Inches(0.3), Inches(0.28), 11, AMB, True)
        txt(sl, desc,  rx + Inches(0.22), ey2 + Inches(0.32), rw - Inches(0.3), Inches(0.78), 10, W2)
        rule(sl, ey2 + Inches(1.18), rx + Inches(0.22), rw - Inches(0.38))
        ey2 += Inches(1.32)

# ── 11 · LIMITATIONS ─────────────────────────────────────────────────
def s11(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "10")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "Limitations & Future Work", "Honest assessment — what works and what to improve next")
    pgnum(sl, 10)
    y0 = content_y()

    lw = (SW - LM - RM - Inches(0.35)) / 2
    rx = LM + lw + Inches(0.35)
    bh = SH - y0 - Inches(0.35)

    # Left: limitations
    box(sl, LM, y0, lw, bh, fill=SURF)
    box(sl, LM, y0, Inches(0.06), bh, fill=AMB)
    txt(sl, "Current Limitations", LM + Inches(0.2), y0 + Inches(0.16),
        lw, Inches(0.3), 13, W1, True)
    rule(sl, y0 + Inches(0.56), LM + Inches(0.2), lw - Inches(0.3))

    lims = [
        ("Recall = 17%",
         "Only one cluster is flagged. Attacks distributed across the remaining four clusters are missed. "
         "Single-cluster flagging is a known trade-off of this approach."),
        ("k must be set manually",
         "The number of clusters is a parameter the user must choose. "
         "In production, the optimal k is not known in advance."),
        ("NSL-KDD is academic",
         "The dataset was created in 1999. It does not represent modern attack types "
         "(ransomware, supply-chain attacks, zero-days)."),
        ("Centroids are random",
         "K-Means results can vary with different random seeds despite n_init=10. "
         "Results are reproducible within this codebase but may differ on other hardware."),
    ]
    lim_h = (bh - Inches(0.66)) / len(lims)
    for i, (title, desc) in enumerate(lims):
        y = y0 + Inches(0.66) + i * lim_h
        txt(sl, title, LM + Inches(0.2), y, lw - Inches(0.3), Inches(0.3), 12, W1, True)
        txt(sl, desc,  LM + Inches(0.2), y + Inches(0.32),
            lw - Inches(0.3), lim_h - Inches(0.38), 10.5, W2)
        if i < len(lims) - 1:
            rule(sl, y + lim_h - Inches(0.06), LM + Inches(0.2), lw - Inches(0.3))

    # Right: future work
    box(sl, rx, y0, lw, bh, fill=SURF)
    box(sl, rx, y0, Inches(0.06), bh, fill=AMB)
    txt(sl, "Next Steps", rx + Inches(0.2), y0 + Inches(0.16),
        lw, Inches(0.3), 13, W1, True)
    rule(sl, y0 + Inches(0.56), rx + Inches(0.2), lw - Inches(0.3))

    futures = [
        ("Auto k selection",
         "Implement Elbow method and Silhouette scan directly in the app. "
         "The user sees a curve and the optimal k is highlighted."),
        ("DBSCAN comparison",
         "Add a second algorithm tab. DBSCAN requires no k and naturally identifies "
         "noise points as anomalies — better recall at cost of interpretability."),
        ("Modern datasets",
         "Replace or supplement NSL-KDD with CICIDS 2017/2018, which includes "
         "modern web attacks, botnet, and DoS vectors."),
        ("Real-time input",
         "Stream live PCAP traffic through the preprocessing pipeline using Scapy. "
         "Classify each new flow against the trained cluster centroids."),
        ("Improve recall",
         "Flag all clusters with attack ratio > threshold, not just the top one. "
         "This trades some precision for significantly higher recall."),
    ]
    fw_h = (bh - Inches(0.66)) / len(futures)
    for i, (title, desc) in enumerate(futures):
        y = y0 + Inches(0.66) + i * fw_h
        box(sl, rx + Inches(0.2), y + Inches(0.1), Inches(0.06), Inches(0.22), fill=AMB)
        txt(sl, title, rx + Inches(0.38), y, lw - Inches(0.5), Inches(0.3), 11.5, W1, True)
        txt(sl, desc,  rx + Inches(0.38), y + Inches(0.32),
            lw - Inches(0.5), fw_h - Inches(0.38), 10.5, W2)
        if i < len(futures) - 1:
            rule(sl, y + fw_h - Inches(0.06), rx + Inches(0.2), lw - Inches(0.3))

# ── 12 · EXPORT ──────────────────────────────────────────────────────
def s12(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "11")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "Export & Analysis History", "Three ways to persist and share results")
    pgnum(sl, 11)
    y0 = content_y()

    cw3 = (SW - LM - RM - 2 * Inches(0.25)) / 3
    bh  = SH - y0 - Inches(0.35)

    panels = [
        ("CSV Download", [
            "Original dataset rows plus two appended columns:",
            "",
            "  cluster        assigned cluster ID (0 to k-1)",
            "  is_anomaly     True / False flag",
            "",
            "Opens in Excel, pandas, R, or any analysis tool.",
            "Ready for downstream ML pipelines.",
        ]),
        ("PDF Report  (fpdf2)", [
            "Auto-generated single-page report containing:",
            "",
            "  File name and algorithm used",
            "  Total records and anomaly count",
            "  Anomaly rate (%)",
            "  Silhouette Score and Davies-Bouldin",
            "  Adjusted Rand Index",
            "  Model parameters for full reproducibility",
        ]),
        ("SQLite History  (database.py)", [
            "Every saved run is stored in analysis_history.db.",
            "",
            "  Columns: timestamp, filename, algorithm,",
            "  n_records, n_anomalies, anomaly_ratio,",
            "  silhouette_score, davies_bouldin_score,",
            "  adjusted_rand_index, params_json",
            "",
            "Auto-migration: schema upgrades gracefully",
            "when new columns are added.",
        ]),
    ]

    for i, (title, lines) in enumerate(panels):
        x = LM + i * (cw3 + Inches(0.25))
        box(sl, x, y0, cw3, bh, fill=SURF)
        box(sl, x, y0, cw3, Inches(0.05), fill=AMB)
        txt(sl, title, x + Inches(0.18), y0 + Inches(0.14),
            cw3 - Inches(0.25), Inches(0.3), 12.5, W1, True)
        rule(sl, y0 + Inches(0.54), x + Inches(0.18), cw3 - Inches(0.3))
        for j, line in enumerate(lines):
            c = W2 if line else W3
            txt(sl, line, x + Inches(0.18),
                y0 + Inches(0.65) + j * Inches(0.42),
                cw3 - Inches(0.28), Inches(0.38), 10.5 if line else 9, c)

# ── 13 · TESTING ─────────────────────────────────────────────────────
def s13(prs):
    sl = blank(prs); bg(sl)
    bignumber_bg(sl, "12")
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    header(sl, "Testing", "44 automated unit tests — pytest across 4 modules, all passing")
    pgnum(sl, 12)
    y0 = content_y()

    # Big 44 banner
    bnh = Inches(1.5)
    box(sl, LM, y0, SW - LM - RM, bnh, fill=SURF)
    box(sl, LM, y0, Inches(0.06), bnh, fill=AMB)
    txt(sl, "44", LM + Inches(0.2), y0 + Inches(0.0), Inches(1.55), Inches(1.18),
        72, AMB, True)
    txt(sl, "unit tests  —  all passing",
        LM + Inches(1.82), y0 + Inches(0.18), Inches(6), Inches(0.38), 20, W1, True)
    txt(sl, "pytest   4 modules   covers loader, preprocessor, model, metrics",
        LM + Inches(1.82), y0 + Inches(0.62), Inches(8), Inches(0.3), 11, W2)
    txt(sl, "Run with:   python -m pytest",
        LM + Inches(1.82), y0 + Inches(0.98), Inches(6), Inches(0.28), 11, W3)

    modules = [
        ("test_loader.py",
         "File loading with CSV / Excel / .txt, column name assignment, schema validation, "
         "error handling for invalid files and missing columns."),
        ("test_preprocessor.py",
         "Output matrix shape after fit-transform, feature name count, correct label separation, "
         "save-to-disk and load-from-disk round-trip for the .pkl file."),
        ("test_models.py",
         "K-Means training on synthetic data, predict returns correct length array, "
         "anomaly mask shape, pre-trained .pkl model loads and predicts without errors."),
        ("test_metrics.py",
         "Silhouette Score, Davies-Bouldin Index, Adjusted Rand Index, "
         "Precision, Recall, F1 Score — all computed correctly on known cluster assignments."),
    ]
    mh = (SH - y0 - bnh - Inches(0.56)) / 4
    for i, (name, desc) in enumerate(modules):
        y = y0 + bnh + Inches(0.18) + i * mh
        fill = SURF if i % 2 == 0 else BG
        box(sl, LM, y + Inches(0.04), SW - LM - RM, mh - Inches(0.06), fill=fill)
        box(sl, LM, y + Inches(0.04), Inches(0.06), mh - Inches(0.06), fill=AMB)
        txt(sl, name, LM + Inches(0.2), y + Inches(0.1),
            SW - LM - RM - Inches(0.28), Inches(0.3), 12, W1, True)
        txt(sl, desc, LM + Inches(0.2), y + Inches(0.42),
            SW - LM - RM - Inches(0.28), mh - Inches(0.5), 10.5, W2)

# ── 14 · CONCLUSION ──────────────────────────────────────────────────
def s14(prs):
    sl = blank(prs); bg(sl)
    box(sl, 0, 0, Inches(0.1), SH, fill=AMB)
    box(sl, 0, SH - Inches(0.44), SW, Inches(0.44), fill=SURF)

    bignumber_bg(sl, "14")
    header(sl, "Conclusion")
    pgnum(sl, 14, 14)

    y0 = content_y(False)

    # 3 metric blocks
    mw = (SW - LM - RM - 2 * Inches(0.14)) / 3
    mh = Inches(1.65)
    for i, (val, lbl, sub) in enumerate([
        ("0.4404",  "Silhouette Score",  "k = 5 — best configuration"),
        ("0.8884",  "Davies-Bouldin",    "Lowest among all tested k"),
        ("44 / 44", "Tests Passing",     "pytest — 4 modules"),
    ]):
        x = LM + i * (mw + Inches(0.14))
        box(sl, x, y0, mw, mh, fill=SURF)
        box(sl, x, y0, mw, Inches(0.05), fill=AMB)
        txt(sl, val, x, y0 + Inches(0.08), mw, Inches(0.9),  38, AMB, True, PP_ALIGN.CENTER)
        txt(sl, lbl, x, y0 + Inches(0.96), mw, Inches(0.3),  11, W1,  True, PP_ALIGN.CENTER)
        txt(sl, sub, x, y0 + Inches(1.25), mw, Inches(0.26), 9.5, W3, False, PP_ALIGN.CENTER)

    rule(sl, y0 + mh + Inches(0.2))
    y2  = y0 + mh + Inches(0.36)
    cw  = (SW - LM - RM - Inches(0.3)) / 2
    rx  = LM + cw + Inches(0.3)
    rbh = SH - y2 - Inches(0.5)

    # Left: delivered
    box(sl, LM, y2, cw, rbh, fill=SURF)
    box(sl, LM, y2, Inches(0.06), rbh, fill=AMB)
    txt(sl, "Delivered", LM + Inches(0.2), y2 + Inches(0.16), cw, Inches(0.3), 13, W1, True)
    items = [
        "Drag-and-drop upload with automatic validation",
        "4 interactive EDA chart tabs (Plotly)",
        "Automated preprocessing — zero configuration",
        "K-Means with adjustable k  +  pre-trained k=5 model",
        "Silhouette, Davies-Bouldin, ARI, Precision, Recall, F1",
        "CSV export, PDF report, SQLite analysis history",
        "44 automated unit tests across 4 pytest modules",
    ]
    ih = (rbh - Inches(0.55)) / len(items)
    for i, item in enumerate(items):
        y = y2 + Inches(0.55) + i * ih
        box(sl, LM + Inches(0.2), y + Inches(0.1), Inches(0.06), Inches(0.22), fill=AMB)
        txt(sl, item, LM + Inches(0.38), y, cw - Inches(0.5), ih, 11, W2)

    # Right: key takeaway
    box(sl, rx, y2, cw, rbh, fill=SURF)
    box(sl, rx, y2, Inches(0.06), rbh, fill=AMB)
    txt(sl, "Key Takeaway", rx + Inches(0.2), y2 + Inches(0.16), cw, Inches(0.3), 13, W1, True)
    txt(sl,
        "Unsupervised clustering can separate normal traffic "
        "from anomalies without any labeled training data — "
        "making it directly deployable in environments where "
        "attack samples are unavailable.\n\n"
        "Precision of 0.9964 confirms the anomaly cluster is "
        "extremely pure. Low recall reflects the single-cluster "
        "constraint — a known limitation addressed honestly "
        "and with concrete improvement paths.\n\n"
        "The complete pipeline runs in a browser with no setup: "
        "upload a file, click two buttons, export the results.",
        rx + Inches(0.2), y2 + Inches(0.56),
        cw - Inches(0.3), rbh - Inches(0.68), 11.5, W1)

    txt(sl, "Thank you  —  Questions welcome",
        LM, SH - Inches(0.36), SW - LM - RM, Inches(0.3),
        12, W3, False, PP_ALIGN.CENTER)

# ── Main ─────────────────────────────────────────────────────────────
def main():
    out = Path(__file__).parent / "PRESENTATION_pro.pptx"
    prs = prs_new()
    s01(prs)   # Title
    s02(prs)   # The Challenge
    s03(prs)   # Dataset
    s04(prs)   # Architecture diagram
    s05(prs)   # 6-step pipeline
    s06(prs)   # Demo: Upload (screenshot)
    s07(prs)   # Demo: EDA (screenshot)
    s08(prs)   # Demo: Results (screenshot)
    s09(prs)   # Demo: Error Analysis (screenshot + why explanation)
    s10(prs)   # Model Selection + why K-Means
    s11(prs)   # Limitations + Future Work
    s12(prs)   # Export & History
    s13(prs)   # Testing
    s14(prs)   # Conclusion
    prs.save(out)
    print(f"Done - {len(prs.slides)} slides saved to {out}")

if __name__ == "__main__":
    main()
